#!/usr/bin/env python3
"""
Konverterar MARP Markdown till PowerPoint (PPTX)

Användning:
    python md_to_pptx.py input.md [output.pptx]

Exempel:
    python md_to_pptx.py slides/lecture_60min.md slides/lecture_60min.pptx
"""

import re
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


class MarpToPptx:
    def __init__(self, dark_theme=True):
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        self.dark_theme = dark_theme

        # Färger baserat på tema
        if dark_theme:
            self.bg_color = RGBColor(26, 26, 26)  # #1a1a1a
            self.text_color = RGBColor(225, 225, 225)  # #e1e1e1
            self.h1_color = RGBColor(74, 158, 255)  # #4a9eff
            self.h2_color = RGBColor(100, 181, 246)  # #64b5f6
            self.code_bg = RGBColor(45, 45, 45)  # #2d2d2d
            self.code_color = RGBColor(168, 224, 99)  # #a8e063
        else:
            self.bg_color = RGBColor(255, 255, 255)
            self.text_color = RGBColor(0, 0, 0)
            self.h1_color = RGBColor(0, 102, 204)
            self.h2_color = RGBColor(51, 122, 183)
            self.code_bg = RGBColor(240, 240, 240)
            self.code_color = RGBColor(51, 51, 51)

    def parse_frontmatter(self, content):
        """Extrahera MARP frontmatter (YAML mellan ---) """
        match = re.match(r'^---\n(.*?)\n---\n', content, re.DOTALL)
        if match:
            frontmatter = match.group(1)
            content = content[match.end():]

            # Kolla om dark theme används
            if 'background-color: #1a1a1a' in frontmatter or 'background: #1a1a1a' in frontmatter:
                self.dark_theme = True
            elif 'background-color: #fff' in frontmatter or 'background: white' in frontmatter:
                self.dark_theme = False

        return content

    def split_slides(self, content):
        """Dela upp markdown i slides baserat på ---"""
        # Ta bort frontmatter först
        content = self.parse_frontmatter(content)

        # Dela vid --- (men inte inom code blocks)
        slides = []
        current_slide = []
        in_code_block = False

        for line in content.split('\n'):
            if line.strip().startswith('```'):
                in_code_block = not in_code_block
                current_slide.append(line)
            elif line.strip() == '---' and not in_code_block:
                if current_slide:
                    slides.append('\n'.join(current_slide))
                    current_slide = []
            else:
                current_slide.append(line)

        # Lägg till sista slide
        if current_slide:
            slides.append('\n'.join(current_slide))

        return slides

    def parse_markdown_line(self, line):
        """Parsa en rad markdown och returnera (text, format_dict)"""
        formats = {
            'is_h1': False,
            'is_h2': False,
            'is_h3': False,
            'is_bold': False,
            'is_italic': False,
            'is_code': False,
            'is_list': False,
            'list_level': 0,
            'is_image_placeholder': False
        }

        # Hantera bilder - konvertera till placeholder
        # Syntax: ![alt text](image_path) eller ![](image_path)
        image_match = re.search(r'!\[([^\]]*)\]\(([^)]+)\)', line)
        if image_match:
            alt_text = image_match.group(1)
            image_path = image_match.group(2)
            # Använd alt text om det finns, annars bildnamnet
            placeholder_text = alt_text if alt_text else image_path.split('/')[-1]
            line = re.sub(r'!\[([^\]]*)\]\(([^)]+)\)', f'[ Plats för bild: {placeholder_text} ]', line)
            formats['is_image_placeholder'] = True

        # Headers
        if line.startswith('# '):
            formats['is_h1'] = True
            line = line[2:]
        elif line.startswith('## '):
            formats['is_h2'] = True
            line = line[3:]
        elif line.startswith('### '):
            formats['is_h3'] = True
            line = line[4:]

        # Lists
        if re.match(r'^(\s*)[-*]\s', line):
            formats['is_list'] = True
            match = re.match(r'^(\s*)[-*]\s', line)
            formats['list_level'] = len(match.group(1)) // 2
            line = re.sub(r'^(\s*)[-*]\s', '', line)

        # Numbered lists
        if re.match(r'^(\s*)\d+\.\s', line):
            formats['is_list'] = True
            match = re.match(r'^(\s*)\d+\.\s', line)
            formats['list_level'] = len(match.group(1)) // 2
            line = re.sub(r'^(\s*)\d+\.\s', '', line)

        # Inline code
        if '`' in line and not line.startswith('```'):
            # Enkel hantering av inline code
            line = line.replace('`', '')
            formats['is_code'] = True

        # Bold (**text** eller __text__)
        if '**' in line or '__' in line:
            formats['is_bold'] = True
            line = re.sub(r'\*\*(.+?)\*\*', r'\1', line)
            line = re.sub(r'__(.+?)__', r'\1', line)

        # Italic (*text* eller _text_)
        if '*' in line or '_' in line:
            formats['is_italic'] = True
            line = re.sub(r'\*(.+?)\*', r'\1', line)
            line = re.sub(r'_(.+?)_', r'\1', line)

        # Ta bort HTML-kommentarer
        line = re.sub(r'<!--.*?-->', '', line)

        # Ta bort emojis från vissa special markers
        line = line.replace('<!-- _class: lead -->', '')

        return line.strip(), formats

    def create_slide(self, content):
        """Skapa en PowerPoint-slide från markdown-innehåll"""
        # Använd blank slide layout
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)

        # Sätt bakgrundsfärg
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.bg_color

        # Parsa innehåll
        lines = content.strip().split('\n')

        # Kolla om det är en title slide (<!-- _class: lead -->)
        is_title_slide = '<!-- _class: lead -->' in content

        if is_title_slide:
            self.create_title_slide(slide, lines)
        else:
            self.create_content_slide(slide, lines)

        return slide

    def create_title_slide(self, slide, lines):
        """Skapa en title slide (centrerad)"""
        # Skapa en stor textbox för hela sliden
        left = Inches(1)
        top = Inches(2)
        width = Inches(14)
        height = Inches(5)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        for line in lines:
            if not line.strip() or '<!--' in line:
                continue

            text, formats = self.parse_markdown_line(line)
            if not text:
                continue

            p = text_frame.add_paragraph()
            p.text = text
            p.alignment = PP_ALIGN.CENTER

            # Formatering
            if formats['is_h1']:
                p.font.size = Pt(54)
                p.font.color.rgb = self.h1_color
                p.font.bold = True
            elif formats['is_h2']:
                p.font.size = Pt(36)
                p.font.color.rgb = self.h2_color
            else:
                p.font.size = Pt(24)
                p.font.color.rgb = self.text_color

    def create_content_slide(self, slide, lines):
        """Skapa en content slide"""
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(15)
        height = Inches(8)

        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        in_code_block = False
        code_lines = []

        for line in lines:
            # Hantera code blocks
            if line.strip().startswith('```'):
                if in_code_block:
                    # Avsluta code block
                    self.add_code_block(text_frame, '\n'.join(code_lines))
                    code_lines = []
                    in_code_block = False
                else:
                    # Starta code block
                    in_code_block = True
                continue

            if in_code_block:
                code_lines.append(line)
                continue

            # Hoppa över tomma rader och kommentarer
            if not line.strip() or '<!--' in line:
                continue

            text, formats = self.parse_markdown_line(line)
            if not text:
                continue

            # Lägg till paragraph
            p = text_frame.add_paragraph()
            p.text = text
            p.level = formats['list_level']

            # Applicera formatering
            if formats['is_h1']:
                p.font.size = Pt(44)
                p.font.color.rgb = self.h1_color
                p.font.bold = True
            elif formats['is_h2']:
                p.font.size = Pt(32)
                p.font.color.rgb = self.h2_color
                p.font.bold = True
            elif formats['is_h3']:
                p.font.size = Pt(28)
                p.font.color.rgb = self.h2_color
            elif formats['is_image_placeholder']:
                # Image placeholder i italic och ljusare färg
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(180, 180, 180)  # Ljusgrå
                p.font.italic = True
            elif formats['is_code']:
                p.font.size = Pt(18)
                p.font.name = 'Consolas'
                p.font.color.rgb = self.code_color
            else:
                p.font.size = Pt(20)
                p.font.color.rgb = self.text_color

            if formats['is_bold']:
                p.font.bold = True
            if formats['is_italic']:
                p.font.italic = True

    def add_code_block(self, text_frame, code):
        """Lägg till en code block med formatering"""
        p = text_frame.add_paragraph()
        p.text = code
        p.font.size = Pt(16)
        p.font.name = 'Consolas'
        p.font.color.rgb = self.code_color
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    def convert(self, md_path, pptx_path=None):
        """Konvertera Markdown till PPTX"""
        md_path = Path(md_path)

        if pptx_path is None:
            pptx_path = md_path.with_suffix('.pptx')
        else:
            pptx_path = Path(pptx_path)

        # Läs markdown
        print(f"📖 Läser {md_path}...")
        content = md_path.read_text(encoding='utf-8')

        # Dela upp i slides
        slides = self.split_slides(content)
        print(f"📊 Hittade {len(slides)} slides")

        # Skapa PPTX-slides
        for i, slide_content in enumerate(slides, 1):
            print(f"   Skapar slide {i}/{len(slides)}...")
            self.create_slide(slide_content)

        # Spara
        print(f"💾 Sparar {pptx_path}...")
        self.prs.save(str(pptx_path))
        print(f"✅ Klar! PowerPoint skapad: {pptx_path}")

        return pptx_path


def main():
    if len(sys.argv) < 2:
        print("Användning: python md_to_pptx.py <input.md> [output.pptx]")
        print()
        print("Exempel:")
        print("  python md_to_pptx.py slides/lecture_60min.md")
        print("  python md_to_pptx.py slides/lecture_60min.md output/presentation.pptx")
        sys.exit(1)

    input_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None

    if not Path(input_file).exists():
        print(f"❌ Fel: Filen '{input_file}' finns inte!")
        sys.exit(1)

    try:
        converter = MarpToPptx(dark_theme=True)
        converter.convert(input_file, output_file)
    except Exception as e:
        print(f"❌ Fel vid konvertering: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
